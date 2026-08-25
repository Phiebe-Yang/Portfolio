#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """創建 PPT 演示文稿"""
    
    # 創建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # 定義顏色
    COLOR_BLUE = RGBColor(31, 78, 121)
    COLOR_LIGHT_BLUE = RGBColor(79, 129, 189)
    COLOR_GREEN = RGBColor(155, 187, 89)
    COLOR_ORANGE = RGBColor(255, 159, 64)
    
    # 第1頁：標題頁
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白佈局
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLUE
    
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "GitHub Copilot CLI\n與文件優化案例"
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "如何將非正式文檔轉換為專業說明書"
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.6))
    date_frame = date_box.text_frame
    date_frame.text = "2026年7月16日"
    date_frame.paragraphs[0].font.size = Pt(16)
    date_frame.paragraphs[0].font.color.rgb = RGBColor(180, 180, 180)
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 第2頁：內容概述
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # 標題
    title2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf2 = title2.text_frame
    tf2.text = "項目目標"
    tf2.paragraphs[0].font.size = Pt(40)
    tf2.paragraphs[0].font.bold = True
    tf2.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 項目列表
    content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    points = [
        ("原始狀況", "RAG 建置說明書是非正式的筆記格式"),
        ("目標", "將其轉換為正式、專業的說明書"),
        ("方法", "通過 GitHub Copilot CLI 與 AI 協作"),
        ("成果", "完成正式 Markdown 和 Word 文檔，並附帶架構圖")
    ]
    
    for i, (label, text) in enumerate(points):
        p = content_frame.add_paragraph()
        p.text = f"{label}："
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_BLUE
        p.level = 0
        p.space_before = Pt(8)
        
        p_content = content_frame.add_paragraph()
        p_content.text = text
        p_content.font.size = Pt(16)
        p_content.level = 1
        p_content.space_before = Pt(2)
    
    # 第3頁：工作流程概覽
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf3 = title3.text_frame
    tf3.text = "工作流程"
    tf3.paragraphs[0].font.size = Pt(40)
    tf3.paragraphs[0].font.bold = True
    tf3.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 流程步驟
    steps = [
        ("1. 問題提出", "用戶詢問如何改進文檔"),
        ("2. 分析現況", "AI 讀取原始 DOCX 文件內容"),
        ("3. 建議方案", "提出正式化改進建議"),
        ("4. 內容重組", "創建結構化 Markdown 說明書"),
        ("5. 格式轉換", "將 Markdown 轉換為 Word 文檔"),
        ("6. 媒體整合", "將架構圖插入文檔")
    ]
    
    content_box = slide3.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, (step, desc) in enumerate(steps):
        p = content_frame.add_paragraph()
        p.text = f"{step}: {desc}"
        p.font.size = Pt(14)
        p.level = 0
        p.space_before = Pt(6)
        
        # 交替顏色
        if i % 2 == 0:
            p.font.color.rgb = COLOR_LIGHT_BLUE
        else:
            p.font.color.rgb = COLOR_GREEN
    
    # 第4頁：溝通方式 - 第一步
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf4 = title4.text_frame
    tf4.text = "第一步：提出需求"
    tf4.paragraphs[0].font.size = Pt(40)
    tf4.paragraphs[0].font.bold = True
    tf4.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 內容
    content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.add_paragraph()
    p.text = "用戶詢問："
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    
    p = content_frame.add_paragraph()
    p.text = '"請幫我做一份ppt內容描述..."'
    p.font.size = Pt(16)
    p.level = 1
    p.font.italic = True
    p.space_before = Pt(6)
    
    p = content_frame.add_paragraph()
    p.text = "AI 回應方式："
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    p.space_before = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "✓ 理解用戶的完整需求"
    p.font.size = Pt(15)
    p.level = 1
    p.space_before = Pt(4)
    
    p = content_frame.add_paragraph()
    p.text = "✓ 確認是否需要對話截圖"
    p.font.size = Pt(15)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "✓ 計劃適當的內容結構"
    p.font.size = Pt(15)
    p.level = 1
    
    # 第5頁：溝通方式 - 第二步
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide5.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf5 = title5.text_frame
    tf5.text = "第二步：查看文件內容"
    tf5.paragraphs[0].font.size = Pt(40)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 內容
    content_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.add_paragraph()
    p.text = "AI 操作："
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    
    p = content_frame.add_paragraph()
    p.text = '1. 使用 view 工具查看資料夾內容'
    p.font.size = Pt(15)
    p.level = 1
    p.space_before = Pt(6)
    
    p = content_frame.add_paragraph()
    p.text = '2. 使用 powershell 解析 DOCX 文件的 XML 內容'
    p.font.size = Pt(15)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = '3. 提取純文本內容進行分析'
    p.font.size = Pt(15)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "提出改進建議："
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    p.space_before = Pt(12)
    
    suggestions = [
        "✓ 添加專業結構（章節、目錄）",
        "✓ 統一語言和術語",
        "✓ 改進排版和格式",
        "✓ 添加表格和圖表"
    ]
    
    for sug in suggestions:
        p = content_frame.add_paragraph()
        p.text = sug
        p.font.size = Pt(14)
        p.level = 1
    
    # 第6頁：溝通方式 - 第三步
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide6.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf6 = title6.text_frame
    tf6.text = "第三步：內容重組與轉換"
    tf6.paragraphs[0].font.size = Pt(40)
    tf6.paragraphs[0].font.bold = True
    tf6.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 內容
    content_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.add_paragraph()
    p.text = "Markdown 文件創建 (RAG建置說明書_正式版.md)："
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_BLUE
    
    p = content_frame.add_paragraph()
    p.text = "• 完整的目錄和章節結構"
    p.font.size = Pt(14)
    p.level = 1
    p.space_before = Pt(4)
    
    p = content_frame.add_paragraph()
    p.text = "• 規範化的格式（代碼塊、表格、列表）"
    p.font.size = Pt(14)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "• 專業的語言和術語"
    p.font.size = Pt(14)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "Word 文檔轉換 (RAG建置說明書.docx)："
    p.font.bold = True
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_BLUE
    p.space_before = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "• 使用 python-docx 進行自動轉換"
    p.font.size = Pt(14)
    p.level = 1
    p.space_before = Pt(4)
    
    p = content_frame.add_paragraph()
    p.text = "• 保留 Markdown 格式（標題、列表、表格等）"
    p.font.size = Pt(14)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "• 檔案大小：約 43.8 KB"
    p.font.size = Pt(14)
    p.level = 1
    
    # 第7頁：溝通方式 - 第四步
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide7.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf7 = title7.text_frame
    tf7.text = "第四步：媒體整合與優化"
    tf7.paragraphs[0].font.size = Pt(40)
    tf7.paragraphs[0].font.bold = True
    tf7.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 內容
    content_box = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.add_paragraph()
    p.text = "用戶要求："
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    
    p = content_frame.add_paragraph()
    p.text = '"請幫我把這張架構圖放進word檔"'
    p.font.size = Pt(15)
    p.level = 1
    p.font.italic = True
    p.space_before = Pt(6)
    
    p = content_frame.add_paragraph()
    p.text = "AI 執行操作："
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_BLUE
    p.space_before = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "1. 複製架構圖到專案目錄"
    p.font.size = Pt(14)
    p.level = 1
    p.space_before = Pt(4)
    
    p = content_frame.add_paragraph()
    p.text = "2. 增強轉換腳本以支持圖片插入"
    p.font.size = Pt(14)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "3. 在附錄部分添加圖片和標題"
    p.font.size = Pt(14)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "4. 重新轉換並生成最終版本（RAG建置說明書_最終版.docx）"
    p.font.size = Pt(14)
    p.level = 1
    
    p = content_frame.add_paragraph()
    p.text = "結果：檔案大小增至 182 KB"
    p.font.bold = True
    p.font.size = Pt(14)
    p.level = 1
    p.space_before = Pt(8)
    p.font.color.rgb = COLOR_GREEN
    
    # 第8頁：AI 工具和技術
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide8.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf8 = title8.text_frame
    tf8.text = "使用的技術工具"
    tf8.paragraphs[0].font.size = Pt(40)
    tf8.paragraphs[0].font.bold = True
    tf8.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 內容
    content_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    tools = [
        ("PowerShell", "讀取檔案、執行命令行操作"),
        ("Python", "DOCX 解析、Markdown 轉換、PPT 生成"),
        ("python-docx", "Word 文檔的創建和編輯"),
        ("python-pptx", "PowerPoint 演示文稿生成"),
        ("Markdown", "結構化文本格式"),
    ]
    
    for tool, desc in tools:
        p = content_frame.add_paragraph()
        p.text = f"{tool}："
        p.font.bold = True
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_LIGHT_BLUE
        p.space_before = Pt(6)
        
        p = content_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(13)
        p.level = 1
    
    # 第9頁：成果展示
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide9.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title9 = slide9.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf9 = title9.text_frame
    tf9.text = "最終成果"
    tf9.paragraphs[0].font.size = Pt(40)
    tf9.paragraphs[0].font.bold = True
    tf9.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 創建表格
    rows, cols = 4, 3
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.4)
    height = Inches(5)
    
    table = slide9.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 設置列寬
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.5)
    
    # 表頭
    headers = ["檔案名稱", "格式", "大小"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        fill = cell.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_BLUE
    
    # 表格內容
    data = [
        ["RAG建置說明書_正式版", "Markdown", "~10 KB"],
        ["RAG建置說明書_最終版", "Word (.docx)", "182 KB"],
        ["RAG_architecture_diagram", "PNG", "149 KB"]
    ]
    
    for i, row in enumerate(data, 1):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = cell_text
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            if i % 2 == 0:
                fill = cell.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    # 第10頁：關鍵收獲
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide10.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    title10 = slide10.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf10 = title10.text_frame
    tf10.text = "關鍵收獲"
    tf10.paragraphs[0].font.size = Pt(40)
    tf10.paragraphs[0].font.bold = True
    tf10.paragraphs[0].font.color.rgb = COLOR_BLUE
    
    # 內容
    content_box = slide10.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(7.6), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    insights = [
        "AI 可以有效地理解並改進非結構化文檔",
        "通過迭代溝通，逐步完善最終成果",
        "自動化工具（Python、PowerShell）大幅提高效率",
        "文檔格式轉換可以無縫進行（MD → DOCX）",
        "媒體整合使文檔更加完整和專業",
        "AI 與人類協作可產出高質量的文檔"
    ]
    
    for i, insight in enumerate(insights):
        p = content_frame.add_paragraph()
        p.text = insight
        p.font.size = Pt(16)
        p.space_before = Pt(8)
        
        # 交替顏色
        if i % 2 == 0:
            p.font.color.rgb = COLOR_LIGHT_BLUE
            p_number = "✓"
        else:
            p.font.color.rgb = COLOR_GREEN
            p_number = "◆"
        
        p.text = p_number + " " + p.text
    
    # 第11頁：結語
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide11.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLUE
    
    # 大標題
    title11 = slide11.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf11 = title11.text_frame
    tf11.text = "感謝使用\nGitHub Copilot CLI"
    tf11.paragraphs[0].font.size = Pt(52)
    tf11.paragraphs[0].font.bold = True
    tf11.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    tf11.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle11 = slide11.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1.5))
    sf11 = subtitle11.text_frame
    sf11.text = "持續改進 • 高效協作 • 優質成果"
    sf11.paragraphs[0].font.size = Pt(28)
    sf11.paragraphs[0].font.color.rgb = RGBColor(200, 220, 255)
    sf11.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 保存演示文稿
    prs.save(r'C:\Users\Phiebe\OneDrive\Desktop\RAG\Copilot_CLI_溝通案例.pptx')
    print('PPT生成成功!')
    print('檔案位置: C:\\Users\\Phiebe\\OneDrive\\Desktop\\RAG\\Copilot_CLI_溝通案例.pptx')

if __name__ == '__main__':
    create_presentation()
